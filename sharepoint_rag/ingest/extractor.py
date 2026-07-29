"""Extract plain text from document bytes for common SharePoint file types.

Heavy parsers (pypdf, python-docx, python-pptx, beautifulsoup4) are imported
lazily so the core package stays light. Install them with the ``extract`` extra.
"""

from __future__ import annotations

import io
from pathlib import Path

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".html",
    ".htm",
    ".md",
    ".markdown",
    ".txt",
    ".csv",
    ".json",
    ".log",
}


class UnsupportedFileType(Exception):
    pass


def is_supported(name: str) -> bool:
    return Path(name).suffix.lower() in SUPPORTED_EXTENSIONS


def extract_text(name: str, data: bytes) -> str:
    """Extract text from ``data`` based on the file extension of ``name``."""
    ext = Path(name).suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(data)
    if ext == ".docx":
        return _extract_docx(data)
    if ext == ".pptx":
        return _extract_pptx(data)
    if ext in (".html", ".htm"):
        return _extract_html(data)
    if ext in (".md", ".markdown", ".txt", ".csv", ".json", ".log"):
        return _decode_text(data)

    raise UnsupportedFileType(f"No extractor for {ext} ({name})")


def _decode_text(data: bytes) -> str:
    # Honor a UTF-16 byte-order mark if present; otherwise UTF-16 would happily
    # decode arbitrary even-length byte strings into mojibake.
    if data[:2] in (b"\xff\xfe", b"\xfe\xff"):
        try:
            return data.decode("utf-16")
        except UnicodeDecodeError:
            pass
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        # latin-1 maps every byte to a code point and never raises.
        return data.decode("latin-1")


def _extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover
        raise ImportError("pypdf required: pip install 'sharepoint-rag-oss[extract]'") from exc

    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_docx(data: bytes) -> str:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover
        raise ImportError("python-docx required: pip install 'sharepoint-rag-oss[extract]'") from exc

    document = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs)


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise ImportError("python-pptx required: pip install 'sharepoint-rag-oss[extract]'") from exc

    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _extract_html(data: bytes) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:  # pragma: no cover
        raise ImportError("beautifulsoup4 required: pip install 'sharepoint-rag-oss[extract]'") from exc

    soup = BeautifulSoup(_decode_text(data), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n")
